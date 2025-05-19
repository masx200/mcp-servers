from typing import Annotated, Optional, Literal, Dict, Any, List, Union
import os
import base64
import io
import sys
import json
from pydantic import BaseModel, Field
import tempfile
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import (
    ErrorData,
    GetPromptResult,
    Prompt,
    PromptArgument,
    PromptMessage,
    TextContent,
    Tool,
    INVALID_PARAMS,
    INTERNAL_ERROR,
)
import pypandoc
import re
import csv
import requests
import yaml
from PIL import Image
from io import BytesIO
import markdown
from bs4 import BeautifulSoup
import PyPDF2
import ghostscript
import camelot
import psutil

# Define custom exception class
class McpError(Exception):
    def __init__(self, error_data: ErrorData):
        self.error_data = error_data
        super().__init__(str(error_data))

# Define content type and insert mode literals
ContentType = Literal["clipboard", "text", "file", "data"]
InsertMode = Literal["append", "overwrite", "merge"]
DocumentType = Literal["doc", "docx", "xls", "xlsx", "ppt", "pptx"]
SaveOption = Literal["true", "false"]

# Error code constants
DOC_LOCKED = "DOC_LOCKED"
FONT_MISSING = "FONT_MISSING"
DATA_OVERFLOW = "DATA_OVERFLOW"
FORMAT_UNSUPPORTED = "FORMAT_UNSUPPORTED"

# Base model for document operations
class DocumentOperationBase(BaseModel):
    target_document: Annotated[str, Field(description="目标文档路径 (doc/docx/xls/xlsx/ppt/pptx)")]
    content_type: Annotated[ContentType, Field(description="插入内容类型 (clipboard/text/file/data)")]
    content: Annotated[str, Field(description="要插入的内容 (根据content_type解析)")]
    position: Annotated[Optional[str], Field(default=None, description="插入位置 (根据文档类型定义不同)")]
    insert_mode: Annotated[Optional[InsertMode], Field(default="append", description="插入模式 (append/overwrite/merge)")]
    image_scale: Annotated[Optional[str], Field(default=None, description="图片缩放 (百分比如50%或固定尺寸如300x200)")]
    table_style: Annotated[Optional[str], Field(default=None, description="表格样式 (如LightGrid/MediumShading)")]
    font_mapping: Annotated[Optional[str], Field(default=None, description="字体映射 (如宋体:SimSun)")]
    overwrite_save: Annotated[Optional[SaveOption], Field(default="false", description="是否覆盖保存 (true/false)")]
    output_format: Annotated[Optional[str], Field(default=None, description="输出格式 (如docx/xlsx/pptx)")]
    compatibility_mode: Annotated[Optional[str], Field(default=None, description="版本兼容 (如Word的97-2003兼容模式)")]
    password: Annotated[Optional[str], Field(default=None, description="加密文档的密码")]

class WriteToDocument(DocumentOperationBase):
    pass

class BatchWriteToDocument(BaseModel):
    input_list: Annotated[str, Field(description="要批量处理的文档列表 (JSON文件路径)")]
    force_unlock: Annotated[Optional[bool], Field(default=False, description="是否强制解除文档占用")]
    auto_expand_rows: Annotated[Optional[bool], Field(default=False, description="是否自动扩展Excel行数")]

class RevertDocument(BaseModel):
    target_document: Annotated[str, Field(description="目标文档路径")]
    version: Annotated[int, Field(description="回退版本号")]

def get_document_type(file_path: str) -> DocumentType:
    """从文件路径获取文档类型"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in ['.doc', '.docx']:
        return 'docx' if ext == '.docx' else 'doc'
    elif ext in ['.xls', '.xlsx']:
        return 'xlsx' if ext == '.xlsx' else 'xls'
    elif ext in ['.ppt', '.pptx']:
        return 'pptx' if ext == '.pptx' else 'ppt'
    else:
        raise ValueError(f"不支持的文档类型: {ext}")

def ensure_dir_exists(file_path: str) -> None:
    """确保文件所在的目录存在"""
    directory = os.path.dirname(file_path)
    if directory and not os.path.exists(directory):
        os.makedirs(directory)

def create_backup(file_path: str) -> str:
    """创建文档备份"""
    if not os.path.exists(file_path):
        return ""
    
    base_name = os.path.splitext(file_path)[0]
    ext = os.path.splitext(file_path)[1]
    
    # Find next version number
    version = 1
    while os.path.exists(f"{base_name}_version{version}.bak"):
        version += 1
    
    backup_path = f"{base_name}_version{version}.bak"
    
    # Copy file to backup
    import shutil
    shutil.copy2(file_path, backup_path)
    
    return backup_path

def get_output_path(target_path: str, overwrite: bool, output_format: Optional[str]) -> str:
    """获取输出文件路径"""
    if overwrite:
        if output_format:
            base_name = os.path.splitext(target_path)[0]
            return f"{base_name}.{output_format}"
        return target_path
    else:
        base_name = os.path.splitext(target_path)[0]
        ext = os.path.splitext(target_path)[1]
        if output_format:
            return f"{base_name}_modified.{output_format}"
        return f"{base_name}_modified{ext}"


def get_content_data(args: WriteToDocument) -> Any:
    """根据content_type获取实际内容数据，增强版"""
    # 常规文本内容
    if args.content_type == "text":
        return args.content
    # 从文件读取
    elif args.content_type == "file":
        file_path = args.content
        if not os.path.exists(file_path):
            raise ValueError(f"指定的文件 {file_path} 不存在")

        # 根据文件扩展名处理不同类型
        file_ext = os.path.splitext(file_path)[1].lower()

        # 文本文件
        if file_ext in ['.txt', '.md', '.html', '.htm', '.json', '.yaml', '.yml', '.csv', '.tsv']:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # 如果是CSV或TSV, 转换为结构化数据
            if file_ext in ['.csv', '.tsv']:
                delimiter = '\t' if file_ext == '.tsv' else ','
                csv_reader = csv.reader(content.splitlines(), delimiter=delimiter)
                return list(csv_reader)

            # 如果是JSON，解析为Python对象
            if file_ext == '.json':
                try:
                    return json.loads(content)
                except:
                    return content  # 如果解析失败，返回原始内容

            # 如果是YAML，解析为Python对象
            if file_ext in ['.yaml', '.yml']:
                try:
                    return yaml.safe_load(content)
                except:
                    return content

            return content

        # PDF文件处理
        elif file_ext == '.pdf':
            # 检查处理模式
            pdf_mode = args.pdf_mode if hasattr(args, 'pdf_mode') else 'image'

            if pdf_mode == 'text':
                # 提取文本内容
                text_content = []
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        text_content.append(page.extract_text())
                return "\n\n".join(text_content)

            elif pdf_mode == 'table':
                # 提取表格内容
                tables = camelot.read_pdf(file_path)
                if len(tables) > 0:
                    result = []
                    for table in tables:
                        result.append(table.df.values.tolist())
                    return result
                else:
                    return "PDF中未检测到表格内容"

            else:  # 默认为图像模式
                # 将PDF页面转换为图像
                try:
                    # 使用临时目录存储图像
                    temp_dir = tempfile.mkdtemp()

                    # 调用GhostScript进行PDF渲染
                    import subprocess
                    gs_cmd = [
                        "gs", "-dNOPAUSE", "-dBATCH", "-dSAFER",
                        "-sDEVICE=pngalpha", "-r300",
                        f"-sOutputFile={temp_dir}/page-%03d.png",
                        file_path
                    ]
                    subprocess.run(gs_cmd, check=True)

                    # 获取生成的图像文件
                    image_files = sorted([f for f in os.listdir(temp_dir) if f.startswith("page-")])

                    # 检查我们是否需要返回单页或多页
                    page_num = args.page_number if hasattr(args, 'page_number') else 1

                    if page_num > 0 and page_num <= len(image_files):
                        # 返回指定页面的图像
                        with open(os.path.join(temp_dir, image_files[page_num - 1]), 'rb') as f:
                            image_data = f.read()

                        # 清理临时文件
                        import shutil
                        shutil.rmtree(temp_dir)

                        return image_data
                    else:
                        # 返回包含所有页面的列表
                        result = []
                        for image_file in image_files:
                            with open(os.path.join(temp_dir, image_file), 'rb') as f:
                                result.append(f.read())

                        # 清理临时文件
                        import shutil
                        shutil.rmtree(temp_dir)

                        return result
                except Exception as e:
                    raise ValueError(f"PDF转图像失败: {str(e)}")

        # 图像文件直接读取二进制数据
        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
            with open(file_path, 'rb') as f:
                image_data = f.read()

            # 使用新函数检查图像分辨率限制
            if not check_image_resolution(image_data):
                raise ValueError("图像分辨率过大，超过10000x10000，拒绝处理")
            return image_data

        # 其他二进制文件直接读取
        else:
            with open(file_path, 'rb') as f:
                return f.read()

    # 从剪贴板获取内容
    elif args.content_type == "clipboard":
        try:
            # 使用新的剪贴板函数
            content = get_clipboard_content()
            if content is None:
                raise ValueError("剪贴板中没有可识别的内容")
            return content
        except Exception as e:
            raise ValueError(f"无法从剪贴板获取内容: {str(e)}")

    # 从数据字段获取
    elif args.content_type == "data":
        # 尝试解析JSON数据
        try:
            # 检查是否为Base64编码
            if re.match(r'^[A-Za-z0-9+/]+={0,2}$', args.content):
                try:
                    # 尝试Base64解码
                    decoded = base64.b64decode(args.content)
                    # 尝试将解码后的内容解析为文本
                    try:
                        text = decoded.decode('utf-8')
                        try:
                            # 尝试解析为JSON
                            return json.loads(text)
                        except:
                            return text
                    except:
                        # 不是UTF-8文本，可能是二进制数据
                        return decoded
                except:
                    # 不是有效的Base64，当作普通文本处理
                    pass

            # 尝试直接解析为JSON
            return json.loads(args.content)
        except json.JSONDecodeError:
            # 尝试解析为YAML
            try:
                return yaml.safe_load(args.content)
            except:
                # 如果都解析失败，检查是否为URL
                if args.content.startswith(('http://', 'https://')):
                    try:
                        # 使用新的HTTP内容获取函数
                        return get_http_content(args.content)
                    except Exception as e:
                        raise ValueError(f"从URL获取数据失败: {str(e)}")

                # 返回原始内容
                return args.content

    else:
        raise ValueError(f"不支持的内容类型: {args.content_type}")

def parse_position_word(position: Optional[str]) -> Dict[str, Any]:
    """解析Word文档的位置参数"""
    if not position:
        return {"position_type": "end"}
    
    # 支持的位置格式:
    # - "start" - 文档开始
    # - "end" - 文档结束
    # - "paragraph:n" - 第n段落后
    # - "bookmark:name" - 指定书签位置
    
    parts = position.split(":")
    if len(parts) == 1:
        if parts[0].lower() in ["start", "end"]:
            return {"position_type": parts[0].lower()}
        else:
            raise ValueError(f"无效的位置参数: {position}")
    
    if len(parts) == 2:
        pos_type, pos_value = parts
        if pos_type.lower() == "paragraph":
            try:
                return {"position_type": "paragraph", "paragraph_index": int(pos_value) - 1}
            except:
                raise ValueError(f"段落索引必须为数字: {pos_value}")
        elif pos_type.lower() == "bookmark":
            return {"position_type": "bookmark", "bookmark_name": pos_value}
    
    raise ValueError(f"无效的位置参数: {position}")

def parse_position_excel(position: Optional[str]) -> Dict[str, Any]:
    """解析Excel文档的位置参数"""
    if not position:
        return {"sheet": "Sheet1", "cell": "A1"}
    
    # 支持的位置格式:
    # - "Sheet1!A1" - 特定工作表和单元格
    # - "A1" - 默认在第一个工作表的单元格
    
    if "!" in position:
        sheet, cell = position.split("!")
        return {"sheet": sheet, "cell": cell}
    else:
        return {"sheet": "Sheet1", "cell": position}

def parse_position_powerpoint(position: Optional[str]) -> Dict[str, Any]:
    """解析PowerPoint文档的位置参数"""
    if not position:
        return {"slide": 0}
    
    # 支持的位置格式:
    # - "1" - 第一张幻灯片 (索引从0开始，所以实际是第0张)
    # - "end" - 最后一张幻灯片
    # - "new" - 创建新幻灯片
    
    if position.lower() == "end":
        return {"slide": "end"}
    elif position.lower() == "new":
        return {"slide": "new"}
    else:
        try:
            return {"slide": int(position) - 1}
        except:
            raise ValueError(f"无效的幻灯片索引: {position}")

def parse_image_scale(scale: Optional[str]) -> Dict[str, Any]:
    """解析图片缩放参数"""
    if not scale:
        return {}
    
    if "%" in scale:
        try:
            percent = float(scale.replace("%", "")) / 100
            return {"percent": percent}
        except:
            raise ValueError(f"无效的百分比缩放值: {scale}")
    elif "x" in scale.lower():
        try:
            width, height = scale.lower().split("x")
            return {"width": int(width), "height": int(height)}
        except:
            raise ValueError(f"无效的尺寸缩放值: {scale}")
    else:
        raise ValueError(f"无效的缩放格式: {scale}")

def parse_font_mapping(mapping: Optional[str]) -> Dict[str, str]:
    """解析字体映射参数"""
    if not mapping:
        return {}
    
    font_map = {}
    for item in mapping.split(";"):
        if ":" in item:
            source, target = item.split(":")
            font_map[source.strip()] = target.strip()
    
    return font_map


def process_excel_document(args: WriteToDocument) -> str:
    """处理Excel文档，增强版"""
    try:
        from openpyxl import load_workbook, Workbook
        from openpyxl.utils import get_column_letter, column_index_from_string
        from openpyxl.styles import PatternFill, Border, Side, Alignment, Font, NamedStyle
        from openpyxl.drawing.image import Image
        import pandas as pd
        import numpy as np
        from datetime import datetime
        import re

        content_data = get_content_data(args)
        position_info = parse_position_excel(args.position)

        # 检查文档是否被占用
        doc_exists = os.path.exists(args.target_document)
        is_temp_file = False

        if doc_exists:
            try:
                # 尝试打开工作簿
                wb = load_workbook(args.target_document)
            except Exception as e:
                error_msg = str(e).lower()
                if "被另一个进程占用" in error_msg or "access denied" in error_msg:
                    # 检查是否允许强制解锁
                    if hasattr(args, 'force_unlock') and args.force_unlock:
                        unlocked = force_unlock_document(args.target_document)
                        if unlocked:
                            # 重新尝试打开
                            try:
                                wb = load_workbook(args.target_document)
                            except:
                                # 强制解锁失败，尝试只读模式
                                raise McpError(ErrorData(code=DOC_LOCKED, message=f"无法解除文档占用: {str(e)}"))
                        else:
                            # 解锁失败，尝试以只读模式打开
                            try:
                                # 创建临时副本
                                temp_path = tempfile.mktemp(suffix='.xlsx')
                                import shutil
                                shutil.copy2(args.target_document, temp_path)
                                wb = load_workbook(temp_path)
                                # 设置标志，表明这是临时文件
                                is_temp_file = True
                            except:
                                raise McpError(
                                    ErrorData(code=DOC_LOCKED, message=f"文档被占用且无法以只读模式打开: {str(e)}"))
                    else:
                        # 未启用强制解锁，直接报错
                        raise McpError(ErrorData(code=DOC_LOCKED, message=f"文档被占用: {str(e)}"))
                elif args.password and ("密码" in str(e) or "password" in error_msg):
                    pass
                    # 尝试使用密码打开（需要win32com）
                    # try:
                    #     # 使用COM接口打开加密文档
                    #     pythoncom.CoInitialize()
                    #     excel_app = win32com.client.Dispatch("Excel.Application")
                    #     excel_app.Visible = False
                    #     excel_app.DisplayAlerts = False
                    #
                    #     # 创建临时文件路径
                    #     temp_path = tempfile.mktemp(suffix='.xlsx')
                    #
                    #     # 打开原始文档并保存为未加密的临时文件
                    #     workbook = excel_app.Workbooks.Open(
                    #         os.path.abspath(args.target_document),
                    #         False,  # 更新链接
                    #         True,  # 只读
                    #         None,  # 格式
                    #         args.password  # 密码
                    #     )
                    #
                    #     workbook.SaveAs(os.path.abspath(temp_path), 51)  # 51 = xlsx format
                    #     workbook.Close()
                    #     excel_app.Quit()
                    #
                    #     # 打开解密后的临时文档
                    #     wb = load_workbook(temp_path)
                    #     # 设置标志，表明这是临时文件
                    #     is_temp_file = True
                    # except Exception as com_error:
                    #     raise McpError(
                    #         ErrorData(code=INVALID_PARAMS, message=f"文档密码错误或无法处理加密文档: {str(com_error)}"))
                else:
                    # 其他错误
                    raise ValueError(f"无法打开文档: {str(e)}")
        else:
            # 新建工作簿
            wb = Workbook()
            # 删除默认的空白工作表
            if 'Sheet' in wb.sheetnames:
                wb.remove(wb['Sheet'])

        # 确定要操作的工作表
        sheet_name = position_info.get("sheet", "Sheet1")

        # 检查位置信息是否包含命名区域
        named_range = None
        if "#" in sheet_name:
            sheet_name, named_range = sheet_name.split("#", 1)

        # 检查工作表是否存在
        if sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
        else:
            ws = wb.create_sheet(title=sheet_name)

        # 处理命名区域（如果有）
        start_cell = position_info.get("cell", "A1")
        if named_range:
            # 查找命名区域
            found = False
            for name in wb.defined_names:
                if name.name == named_range:
                    found = True
                    # 获取范围值
                    destinations = list(name.destinations)
                    if destinations:
                        sheet, coord = destinations[0]
                        if sheet == sheet_name:
                            start_cell = coord.split(':')[0]  # 使用命名区域的左上角单元格
                    break

            if not found:
                # 命名区域不存在，使用默认单元格
                pass

        # 解析单元格位置
        col, row = None, None

        # 解析单元格引用 (如A1)
        cell_match = re.match(r"([A-Za-z]+)(\d+)", start_cell)
        if cell_match:
            col_str, row_str = cell_match.groups()
            col = column_index_from_string(col_str)
            row = int(row_str)
        else:
            # 默认位置
            col = 1
            row = 1

        # 获取插入模式
        insert_mode = args.insert_mode or "append"

        # 根据内容类型处理
        if isinstance(content_data, bytes):  # 图片数据
            # 创建临时文件保存图片
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                tmp.write(content_data)
                tmp_name = tmp.name

            # 插入图片
            img = Image(tmp_name)

            # 应用缩放
            image_scale = parse_image_scale(args.image_scale)
            if "percent" in image_scale:
                img.width = int(img.width * image_scale["percent"])
                img.height = int(img.height * image_scale["percent"])
            elif "width" in image_scale and "height" in image_scale:
                img.width = image_scale["width"]
                img.height = image_scale["height"]

            # 添加图片到工作表
            cell_address = f"{get_column_letter(col)}{row}"
            ws.add_image(img, cell_address)

            # 删除临时文件
            os.unlink(tmp_name)

        elif isinstance(content_data, list):  # 列表/表格数据
            # 处理二维数组数据
            if all(isinstance(row, list) for row in content_data):
                # 检查是否需要自动展开行
                auto_expand = hasattr(args, 'auto_expand_rows') and args.auto_expand_rows

                # 准备写入数据前，检查数据量
                data_rows = len(content_data)
                data_cols = max(len(row) for row in content_data) if data_rows > 0 else 0

                # 如果数据过大，启用流式处理
                if data_rows * data_cols > 100000:  # 假设超过10万单元格就是大数据量
                    # 提示将使用流式处理
                    print(f"检测到大数据量 ({data_rows}行 x {data_cols}列)，启用流式处理...")

                    # 使用分块写入
                    chunk_size = 5000  # 每次处理的行数
                    for chunk_start in range(0, data_rows, chunk_size):
                        chunk_end = min(chunk_start + chunk_size, data_rows)
                        chunk_data = content_data[chunk_start:chunk_end]

                        for i, row_data in enumerate(chunk_data):
                            row_index = row + chunk_start + i
                            for j, cell_value in enumerate(row_data):
                                if j < data_cols:  # 确保不超出列范围
                                    cell = ws.cell(row=row_index, column=col + j)

                                    # 如果是覆盖模式或单元格为空，则写入
                                    if insert_mode == "overwrite" or cell.value is None:
                                        cell.value = cell_value
                                    # 如果是合并模式且单元格有值，尝试合并
                                    elif insert_mode == "merge":
                                        # 如果单元格包含公式，保留公式
                                        if isinstance(cell.value, str) and cell.value.startswith("="):
                                            pass  # 保持原公式不变
                                        else:
                                            # 合并数值 - 对于数字，尝试相加；对于文本，尝试连接
                                            if isinstance(cell.value, (int, float)) and isinstance(cell_value,
                                                                                                   (int, float)):
                                                cell.value += cell_value
                                            else:
                                                cell.value = f"{cell.value}, {cell_value}"

                        # 每批次写入后刷新
                        wb.save(args.target_document + "_temp")

                    # 处理完成后，复制临时文件到目标位置
                    os.replace(args.target_document + "_temp", args.target_document)
                else:
                    # 常规处理 - 数据量不大
                    for i, row_data in enumerate(content_data):
                        for j, cell_value in enumerate(row_data):
                            cell = ws.cell(row=row + i, column=col + j)

                            # 根据插入模式处理
                            if insert_mode == "overwrite" or cell.value is None:
                                cell.value = cell_value
                            elif insert_mode == "merge":
                                # 如果是公式，保留原公式
                                if isinstance(cell.value, str) and cell.value.startswith("="):
                                    pass
                                # 如果是数字，尝试相加
                                elif isinstance(cell.value, (int, float)) and isinstance(cell_value, (int, float)):
                                    cell.value += cell_value
                                # 否则连接字符串
                                else:
                                    cell.value = f"{cell.value}, {cell_value}"

                # 检测和应用数据格式
                auto_format_data(ws, row, col, content_data, data_rows, data_cols)

                # 应用表格样式
                if args.table_style:
                    apply_excel_style(ws, row, col, data_rows, data_cols, args.table_style)

            # 单维列表 - 作为单行或单列处理
            else:
                # 判断放置方向 - 默认横向(行)
                direction = args.direction if hasattr(args, 'direction') else "row"

                if direction.lower() == "column":
                    # 纵向放置(列)
                    for i, item in enumerate(content_data):
                        cell = ws.cell(row=row + i, column=col)
                        if insert_mode == "overwrite" or cell.value is None:
                            cell.value = item
                        elif insert_mode == "merge":
                            if isinstance(cell.value, (int, float)) and isinstance(item, (int, float)):
                                cell.value += item
                            else:
                                cell.value = f"{cell.value}, {item}"
                else:
                    # 横向放置(行)
                    for j, item in enumerate(content_data):
                        cell = ws.cell(row=row, column=col + j)
                        if insert_mode == "overwrite" or cell.value is None:
                            cell.value = item
                        elif insert_mode == "merge":
                            if isinstance(cell.value, (int, float)) and isinstance(item, (int, float)):
                                cell.value += item
                            else:
                                cell.value = f"{cell.value}, {item}"

        elif isinstance(content_data, dict):  # 字典数据
            # 转换为两列表格：键和值
            keys = list(content_data.keys())
            values = list(content_data.values())

            # 默认纵向放置: 第一列是键，第二列是值
            for i, (key, value) in enumerate(zip(keys, values)):
                # 键
                key_cell = ws.cell(row=row + i, column=col)
                if insert_mode == "overwrite" or key_cell.value is None:
                    key_cell.value = key

                # 值
                value_cell = ws.cell(row=row + i, column=col + 1)
                if insert_mode == "overwrite" or value_cell.value is None:
                    value_cell.value = value
                elif insert_mode == "merge":
                    if isinstance(value_cell.value, (int, float)) and isinstance(value, (int, float)):
                        value_cell.value += value
                    else:
                        value_cell.value = f"{value_cell.value}, {value}"

            # 应用字典样式
            if args.table_style:
                apply_excel_style(ws, row, col, len(keys), 2, args.table_style)

        else:  # 字符串或其他类型内容
            content_str = str(content_data)

            # 检查是否为CSV文本
            if "\n" in content_str and ("," in content_str or "\t" in content_str):
                # 尝试解析为CSV
                delimiter = "\t" if "\t" in content_str else ","

                try:
                    # 使用pandas解析
                    with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.csv') as tmp:
                        tmp.write(content_str)
                        tmp_name = tmp.name

                    # 尝试自动识别日期和数字格式
                    df = pd.read_csv(tmp_name, delimiter=delimiter, parse_dates=True, infer_datetime_format=True)

                    # 获取数据维度
                    data_rows, data_cols = df.shape

                    # 写入表头
                    for j, col_name in enumerate(df.columns):
                        ws.cell(row=row, column=col + j, value=col_name)

                    # 写入数据
                    for i in range(data_rows):
                        for j in range(data_cols):
                            cell_value = df.iloc[i, j]

                            # 处理NaN值
                            if isinstance(cell_value, float) and np.isnan(cell_value):
                                cell_value = None

                            # 处理日期值
                            if isinstance(cell_value, pd.Timestamp):
                                cell_value = cell_value.to_pydatetime()

                            cell = ws.cell(row=row + i + 1, column=col + j)

                            if insert_mode == "overwrite" or cell.value is None:
                                cell.value = cell_value
                            elif insert_mode == "merge":
                                if isinstance(cell.value, (int, float)) and isinstance(cell_value, (int, float)):
                                    cell.value += cell_value
                                else:
                                    cell.value = f"{cell.value}, {cell_value}"

                    # 应用表格样式
                    if args.table_style:
                        apply_excel_style(ws, row, col, data_rows + 1, data_cols, args.table_style)

                    # 删除临时文件
                    os.unlink(tmp_name)
                except Exception as e:
                    # CSV解析失败，作为普通文本添加
                    cell = ws.cell(row=row, column=col)
                    if insert_mode == "overwrite" or cell.value is None:
                        cell.value = content_str
                    elif insert_mode == "merge":
                        cell.value = f"{cell.value}\n{content_str}"
            else:
                # 普通文本，直接添加到单元格
                cell = ws.cell(row=row, column=col)
                if insert_mode == "overwrite" or cell.value is None:
                    cell.value = content_str
                elif insert_mode == "merge":
                    if isinstance(cell.value, (int, float)) and content_str.replace('.', '', 1).isdigit():
                        # 如果原值是数字且新值也可以转为数字
                        try:
                            numeric_value = float(content_str)
                            cell.value += numeric_value
                        except:
                            cell.value = f"{cell.value}\n{content_str}"
                    else:
                        cell.value = f"{cell.value}\n{content_str}"

        # 保存工作簿
        output_path = get_output_path(args.target_document, args.overwrite_save == 'true', args.output_format)

        # 确保输出目录存在
        ensure_dir_exists(output_path)

        wb.save(output_path)

        # 删除临时文件
        if is_temp_file and os.path.exists(temp_path):
            os.unlink(temp_path)

        return f"Excel文档处理成功，内容已插入到 {output_path}"

    except ImportError as e:
        if "openpyxl" in str(e):
            raise ImportError("openpyxl")
        elif "pandas" in str(e):
            raise ImportError("pandas")
        elif "numpy" in str(e):
            raise ImportError("numpy")
        elif "com" in str(e) or "win32" in str(e):
            raise ImportError("pywin32")
        else:
            raise ImportError(str(e))
    except Exception as e:
        if "单元格溢出" in str(e) or "数据超出范围" in str(e):
            raise McpError(ErrorData(code=DATA_OVERFLOW, message=str(e)))
        raise Exception(f"处理Excel文档时出错: {str(e)}")


def auto_format_data(ws, start_row, start_col, data, rows, cols):
    """自动检测和应用数据格式"""
    from openpyxl.styles import numbers

    # 检查每列数据类型
    for j in range(cols):
        # 跳过首行（可能是标题）
        column_data = [row[j] for row in data[1:] if j < len(row)]

        if not column_data:
            continue

        # 检查是否为日期列
        date_count = 0
        for value in column_data:
            if isinstance(value, str):
                # 检查常见日期格式
                date_patterns = [
                    r'\d{4}-\d{1,2}-\d{1,2}',  # YYYY-MM-DD
                    r'\d{1,2}/\d{1,2}/\d{4}',  # MM/DD/YYYY
                    r'\d{1,2}-\d{1,2}-\d{4}',  # DD-MM-YYYY
                    r'\d{4}/\d{1,2}/\d{1,2}',  # YYYY/MM/DD
                    r'\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\s+\d{4}'  # DD MMM YYYY
                ]

                if any(re.match(pattern, value) for pattern in date_patterns):
                    date_count += 1

        # 如果超过50%的值看起来像日期
        if date_count > len(column_data) * 0.5:
            # 应用日期格式
            for i in range(rows):
                cell = ws.cell(row=start_row + i, column=start_col + j)
                cell.number_format = numbers.FORMAT_DATE_YYYYMMDD
            continue

        # 检查是否为货币列
        currency_count = 0
        for value in column_data:
            if isinstance(value, str):
                # 检查常见货币格式
                currency_patterns = [
                    r'^\$\s*\d+(\.\d+)?$',  # $123.45
                    r'^\d+(\.\d+)?\s*\$$',  # 123.45$
                    r'^\d+(\.\d+)?\s*元$',  # 123.45元
                    r'^\d+(\.\d+)?\s*€$',  # 123.45€
                    r'^\€\s*\d+(\.\d+)?$',  # €123.45
                    r'^\£\s*\d+(\.\d+)?$',  # £123.45
                    r'^\d+(\.\d+)?\s*\£$',  # 123.45£
                ]

                if any(re.match(pattern, value) for pattern in currency_patterns):
                    currency_count += 1

            # 如果超过50%的值看起来像货币
        if currency_count > len(column_data) * 0.5:
            # 应用货币格式
            for i in range(rows):
                cell = ws.cell(row=start_row + i, column=start_col + j)
                # 使用基本货币格式
                cell.number_format = numbers.FORMAT_CURRENCY_USD_SIMPLE
            continue

            # 检查是否为百分比列
        percent_count = 0
        for value in column_data:
            if isinstance(value, str) and (value.endswith('%') or value.endswith('％')):
                percent_count += 1
            elif isinstance(value, (int, float)) and 0 <= value <= 1:
                # 可能是小数形式的百分比（0-1之间）
                percent_count += 1

        # 如果超过50%的值看起来像百分比
        if percent_count > len(column_data) * 0.5:
            # 应用百分比格式
            for i in range(rows):
                cell = ws.cell(row=start_row + i, column=start_col + j)
                cell.number_format = numbers.FORMAT_PERCENTAGE
            continue

        # 检查是否为数字列
        number_count = 0
        for value in column_data:
            if isinstance(value, (int, float)):
                number_count += 1
            elif isinstance(value, str) and re.match(r'^-?\d+(\.\d+)?$', value):
                number_count += 1

        # 如果超过70%的值看起来像数字
        if number_count > len(column_data) * 0.7:
            # 应用常规数字格式
            for i in range(rows):
                cell = ws.cell(row=start_row + i, column=start_col + j)
                # 决定是否使用小数位
                if any(isinstance(val, float) for val in column_data):
                    cell.number_format = numbers.FORMAT_NUMBER_00
                else:
                    cell.number_format = numbers.FORMAT_NUMBER

def apply_excel_style(ws, start_row, start_col, rows, cols, style_name):
    """应用Excel表格样式"""
    from openpyxl.styles import PatternFill, Border, Side, Alignment, Font, NamedStyle

    # 定义预设样式
    styles = {
        "LightGrid": {
            "header": {
                "fill": PatternFill(start_color="D9D9D9", end_color="D9D9D9", fill_type="solid"),
                "font": Font(bold=True),
                "border": Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                ),
                "alignment": Alignment(horizontal="center", vertical="center")
            },
            "body": {
                "border": Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                ),
                "alignment": Alignment(vertical="center")
            },
            "alt_row": PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
        },
        "MediumShading": {
            "header": {
                "fill": PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid"),
                "font": Font(bold=True, color="FFFFFF"),
                "border": Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                ),
                "alignment": Alignment(horizontal="center", vertical="center")
            },
            "body": {
                "border": Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                ),
                "alignment": Alignment(vertical="center")
            },
            "alt_row": PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
        },
        "Accent1": {
            "header": {
                "fill": PatternFill(start_color="5B9BD5", end_color="5B9BD5", fill_type="solid"),
                "font": Font(bold=True, color="FFFFFF"),
                "border": Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                ),
                "alignment": Alignment(horizontal="center", vertical="center")
            },
            "body": {
                "border": Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                ),
                "alignment": Alignment(vertical="center")
            },
            "alt_row": PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
        }
    }

    # 默认使用LightGrid样式
    style = styles.get(style_name, styles["LightGrid"])

    # 应用表头样式（第一行）
    for j in range(cols):
        cell = ws.cell(row=start_row, column=start_col + j)
        cell.fill = style["header"]["fill"]
        cell.font = style["header"]["font"]
        cell.border = style["header"]["border"]
        cell.alignment = style["header"]["alignment"]

    # 应用表体样式（从第二行开始）
    for i in range(1, rows):
        for j in range(cols):
            cell = ws.cell(row=start_row + i, column=start_col + j)
            cell.border = style["body"]["border"]
            cell.alignment = style["body"]["alignment"]

            # 应用交替行颜色
            if i % 2 == 1:  # 奇数行（从0开始计数）
                cell.fill = style.get("alt_row", PatternFill())


def force_unlock_document(file_path: str) -> bool:
    """尝试强制解锁被占用的文档（Windows系统专用）"""
    try:
        if os.name != 'nt':  # 只在Windows系统下尝试
            return False

        # 获取所有进程
        occupied = False
        target_pid = None

        for proc in psutil.process_iter(['pid', 'name', 'open_files']):
            try:
                # 检查进程是否有打开的文件
                if proc.info['open_files']:
                    for file in proc.info['open_files']:
                        if file.path == os.path.abspath(file_path):
                            occupied = True
                            target_pid = proc.info['pid']
                            proc_name = proc.info['name']
                            print(f"文件 {file_path} 被进程 {proc_name} (PID: {target_pid}) 占用")
                            break
                    if occupied:
                        break
            except (psutil.AccessDenied, psutil.NoSuchProcess):
                # 无法访问某些进程信息，跳过
                continue

        if not occupied or not target_pid:
            # 未找到占用进程
            return False

        # 尝试终止占用进程
        try:
            process = psutil.Process(target_pid)
            process_name = process.name()

            # 检查是否为Office进程
            office_procs = ['WINWORD.EXE', 'EXCEL.EXE', 'POWERPNT.EXE', 'OUTLOOK.EXE']
            if process_name.upper() in office_procs:
                # 是Office进程，尝试优雅关闭
                print(f"尝试关闭Office进程: {process_name} (PID: {target_pid})")

                # if process_name.upper() == 'WINWORD.EXE':
                #     # 使用COM接口关闭Word
                #     try:
                #         pythoncom.CoInitialize()
                #         word_app = win32com.client.GetObject("Word.Application")
                #         # 检查是否有打开的文档
                #         if word_app.Documents.Count > 0:
                #             # 尝试保存并关闭文档
                #             for i in range(word_app.Documents.Count):
                #                 doc = word_app.Documents.Item(i + 1)
                #                 if os.path.abspath(doc.FullName) == os.path.abspath(file_path):
                #                     # 找到目标文档，保存并关闭
                #                     doc.Close(SaveChanges=True)
                #                     print(f"成功关闭Word文档")
                #                     return True
                #     except:
                #         # COM接口失败，尝试终止进程
                #         pass

                # elif process_name.upper() == 'EXCEL.EXE':
                #     # 使用COM接口关闭Excel
                #     try:
                #         # pythoncom.CoInitialize()
                #         # excel_app = win32com.client.GetObject("Excel.Application")
                #         # # 检查是否有打开的工作簿
                #         # if excel_app.Workbooks.Count > 0:
                #         #     # 尝试保存并关闭工作簿
                #         #     for i in range(excel_app.Workbooks.Count):
                #         #         wb = excel_app.Workbooks.Item(i + 1)
                #         #         if os.path.abspath(wb.FullName) == os.path.abspath(file_path):
                #         #             # 找到目标工作簿，保存并关闭
                #         #             wb.Close(SaveChanges=True)
                #         #             print(f"成功关闭Excel工作簿")
                #                     return True
                #     except:
                #         # COM接口失败，尝试终止进程
                #         pass
                #
                # elif process_name.upper() == 'POWERPNT.EXE':
                #     # 使用COM接口关闭PowerPoint
                #     # try:
                #         # pythoncom.CoInitialize()
                #         # ppt_app = win32com.client.GetObject("PowerPoint.Application")
                #         # # 检查是否有打开的演示文稿
                #         # if ppt_app.Presentations.Count > 0:
                #         #     # 尝试保存并关闭演示文稿
                #         #     for i in range(ppt_app.Presentations.Count):
                #         #         pres = ppt_app.Presentations.Item(i + 1)
                #         #         if os.path.abspath(pres.FullName) == os.path.abspath(file_path):
                #         #             # 找到目标演示文稿，保存并关闭
                #         #             pres.Close()
                #         #             print(f"成功关闭PowerPoint演示文稿")
                #         #             return True
                #     except:
                #         # COM接口失败，尝试终止进程
                #         pass

            # 如果之前的方法未成功，直接终止进程
            print(f"尝试终止进程: {process_name} (PID: {target_pid})")
            process.terminate()
            process.wait(5)  # 等待进程终止
            print(f"进程已终止")
            return True

        except Exception as e:
            print(f"解锁文档时出错: {str(e)}")
            return False

    except Exception as e:
        print(f"强制解锁文档失败: {str(e)}")
        return False


def parse_position_word(position: str) -> dict:
    """解析Word文档的位置信息"""
    position_info = {"position_type": "end", "paragraph_index": 0, "bookmark_name": ""}

    if not position:
        return position_info

    # 处理书签
    if position.startswith("Bookmark") or ":" not in position:
        position_info["position_type"] = "bookmark"
        position_info["bookmark_name"] = position
    # 处理段落索引
    elif position.startswith("p:"):
        position_info["position_type"] = "paragraph"
        try:
            position_info["paragraph_index"] = int(position[2:])
        except ValueError:
            position_info["paragraph_index"] = 0
    # 处理开头/结尾
    elif position.lower() == "start":
        position_info["position_type"] = "start"
    elif position.lower() == "end":
        position_info["position_type"] = "end"

    return position_info


def parse_position_excel(position: str) -> dict:
    """解析Excel工作表的位置信息"""
    position_info = {"sheet": "Sheet1", "cell": "A1"}

    if not position:
        return position_info

    # 处理工作表!单元格格式
    if "!" in position:
        sheet, cell = position.split("!", 1)
        position_info["sheet"] = sheet
        position_info["cell"] = cell
    # 处理命名区域
    elif "#" in position or ":" not in position:
        # 可能是命名区域或仅工作表名
        position_info["sheet"] = position
    # 处理单元格引用
    elif re.match(r"[A-Za-z]+\d+", position):
        position_info["cell"] = position

    return position_info


def parse_position_powerpoint(position: str) -> dict:
    """解析PowerPoint幻灯片的位置信息"""
    position_info = {"slide": "end", "placeholder": ""}

    if not position:
        return position_info

    # 处理slide:n格式
    if position.startswith("slide:"):
        try:
            # 检查是否有占位符指定
            if "#" in position:
                slide_part, placeholder = position.split("#", 1)
                position_info["slide"] = int(slide_part[6:])
                position_info["placeholder"] = placeholder
            else:
                position_info["slide"] = int(position[6:])
        except ValueError:
            # 无效的数字，使用默认值"end"
            pass
    # 处理'new'或'end'
    elif position == "new":
        position_info["slide"] = "new"
    elif position == "end":
        position_info["slide"] = "end"
    # 处理纯数字（幻灯片索引）
    elif position.isdigit():
        position_info["slide"] = int(position)

    return position_info


def parse_image_scale(scale_str: str) -> dict:
    """解析图片缩放参数"""
    if not scale_str:
        return {}

    # 百分比格式，如 "50%"
    if scale_str.endswith("%"):
        try:
            percent = float(scale_str.rstrip("%")) / 100.0
            return {"percent": percent}
        except ValueError:
            return {}

    # 固定尺寸格式，如 "300x200"
    if "x" in scale_str:
        try:
            width, height = scale_str.split("x", 1)
            return {"width": int(width), "height": int(height)}
        except ValueError:
            return {}

    return {}


def parse_font_mapping(mapping_str: str) -> dict:
    """解析字体映射参数"""
    if not mapping_str:
        return {}

    font_map = {}
    try:
        # 处理多个映射，格式如 "宋体:SimSun,黑体:SimHei"
        for pair in mapping_str.split(","):
            if ":" in pair:
                src, dst = pair.split(":", 1)
                font_map[src.strip()] = dst.strip()
    except:
        # 解析失败，返回空映射
        pass

    return font_map


def get_output_path(target_document: str, overwrite: bool, output_format: str = None) -> str:
    """获取输出文件路径"""
    base_name, ext = os.path.splitext(target_document)

    # 如果指定了输出格式，更改扩展名
    if output_format:
        if not output_format.startswith("."):
            output_format = "." + output_format
        ext = output_format

    # 决定是否覆盖原文件
    if overwrite:
        # 先创建备份
        create_backup(target_document)
        return target_document
    else:
        # 创建新文件，添加_modified后缀
        return f"{base_name}_modified{ext}"


def create_backup(file_path: str) -> str:
    """创建文档备份"""
    if not os.path.exists(file_path):
        return None

    # 查找现有备份
    base_name = os.path.splitext(file_path)[0]
    ext = os.path.splitext(file_path)[1]

    # 找到最高版本号
    version = 1
    while os.path.exists(f"{base_name}_version{version}.bak"):
        version += 1

    # 创建新备份
    backup_path = f"{base_name}_version{version}.bak"
    import shutil
    shutil.copy2(file_path, backup_path)

    return backup_path


def ensure_dir_exists(file_path: str) -> None:
    """确保目录存在"""
    dir_name = os.path.dirname(file_path)
    if dir_name and not os.path.exists(dir_name):
        os.makedirs(dir_name)

# 错误代码常量
DOC_LOCKED = "DOC_LOCKED"
FONT_MISSING = "FONT_MISSING"
DATA_OVERFLOW = "DATA_OVERFLOW"
INVALID_PARAMS = "INVALID_PARAMS"
INTERNAL_ERROR = "INTERNAL_ERROR"

# 自定义错误类
class McpError(Exception):
    def __init__(self, error_data: ErrorData):
        self.error_data = error_data
        super().__init__(error_data.message)


async def process_document(args: WriteToDocument) -> str:
    """处理文档写入操作"""
    try:
        # 确认是否支持的文档类型
        doc_type = get_document_type(args.target_document)
        
        # 确保目录存在
        ensure_dir_exists(args.target_document)
        
        # 如果文档不存在且不是data内容，或者是doc_type为data但未指定保存格式，则报错
        if not os.path.exists(args.target_document) and args.content_type != "data":
            raise ValueError(f"目标文档 {args.target_document} 不存在")
        
        # 创建备份（如果文档存在）
        backup_path = ""
        if os.path.exists(args.target_document):
            backup_path = create_backup(args.target_document)
        
        # 处理不同类型的文档
        result = ""
        
        # 检查必要的依赖库
        try:
            if doc_type in ['doc', 'docx']:
                result = process_word_document(args)
            elif doc_type in ['xls', 'xlsx']:
                result = process_excel_document(args)
            elif doc_type in ['ppt', 'pptx']:
                result = process_powerpoint_document(args)
        except ImportError as e:
            missing_lib = str(e).split("'")[1]
            return f"错误：缺少必要的库 {missing_lib}，请安装: pip install {missing_lib}"
        
        if backup_path:
            result = f"{result}\n(已创建备份: {backup_path})"
        
        output_path = get_output_path(args.target_document, args.overwrite_save == 'true', args.output_format)
        return result if result else f"文档 {args.target_document} 处理完成，输出位置: {output_path}"
    
    except Exception as e:
        if "被另一个进程占用" in str(e) or "access denied" in str(e).lower():
            raise McpError(ErrorData(code=DOC_LOCKED, message=f"文档被占用: {str(e)}"))
        if "缺少字体" in str(e) or "font not found" in str(e).lower():
            raise McpError(ErrorData(code=FONT_MISSING, message=f"缺失字体: {str(e)}"))
        if "数据超出范围" in str(e) or "cell overflow" in str(e).lower():
            raise McpError(ErrorData(code=DATA_OVERFLOW, message=f"Excel单元格溢出: {str(e)}"))
        if "格式不支持" in str(e) or "format not supported" in str(e).lower():
            raise McpError(ErrorData(code=FORMAT_UNSUPPORTED, message=f"不支持的格式: {str(e)}"))
        raise McpError(ErrorData(code=INTERNAL_ERROR, message=f"处理文档时出错: {str(e)}"))


def markdown_to_docx(markdown_text: str, doc) -> None:
    """将Markdown文本转换为Word文档格式"""
    try:
        # 首先尝试使用pypandoc (需要安装pandoc)
        try:
            import pypandoc

            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.md', delete=False) as tmp_md:
                tmp_md.write(markdown_text.encode('utf-8'))
                tmp_md_path = tmp_md.name

            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_docx:
                tmp_docx_path = tmp_docx.name

            # 使用pandoc转换
            pypandoc.convert_file(tmp_md_path, 'docx', outputfile=tmp_docx_path)

            # 读取转换后的docx
            from docx import Document
            converted_doc = Document(tmp_docx_path)

            # 将转换后的内容添加到原始文档
            for element in converted_doc.element.body:
                doc.element.body.append(element)

            # 清理临时文件
            os.unlink(tmp_md_path)
            os.unlink(tmp_docx_path)

            return
        except (ImportError, Exception):
            # 如果pypandoc不可用或失败，使用替代方法
            pass

        # 备选方法：使用markdown和python-docx手动转换
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

        # 将Markdown转换为HTML
        html = markdown.markdown(markdown_text, extensions=['tables', 'fenced_code'])

        # 解析HTML
        soup = BeautifulSoup(html, 'html.parser')

        # 递归处理HTML元素
        def process_element(element, parent):
            if element.name is None:
                # 文本节点
                if element.string and element.string.strip():
                    if isinstance(parent, type(doc)):
                        p = parent.add_paragraph()
                        p.add_run(element.string)
                    else:
                        parent.add_run(element.string)
            elif element.name == 'p':
                p = doc.add_paragraph()
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'h1':
                p = doc.add_heading(level=1)
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'h2':
                p = doc.add_heading(level=2)
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'h3':
                p = doc.add_heading(level=3)
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'h4':
                p = doc.add_heading(level=4)
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'h5':
                p = doc.add_heading(level=5)
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'h6':
                p = doc.add_heading(level=6)
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'ul':
                for li in element.find_all('li', recursive=False):
                    p = doc.add_paragraph(style='List Bullet')
                    for child in li.children:
                        process_element(child, p)
            elif element.name == 'ol':
                for i, li in enumerate(element.find_all('li', recursive=False)):
                    p = doc.add_paragraph(style='List Number')
                    for child in li.children:
                        process_element(child, p)
            elif element.name == 'strong' or element.name == 'b':
                run = parent.add_run(element.get_text())
                run.bold = True
            elif element.name == 'em' or element.name == 'i':
                run = parent.add_run(element.get_text())
                run.italic = True
            elif element.name == 'a':
                run = parent.add_run(element.get_text())
                run.underline = True
                run.font.color.rgb = RGBColor(0, 0, 255)
                if 'href' in element.attrs:
                    run.hyperlink = element['href']
            elif element.name == 'code':
                run = parent.add_run(element.get_text())
                run.font.name = 'Courier New'
            elif element.name == 'pre':
                p = doc.add_paragraph()
                code_text = element.get_text()
                run = p.add_run(code_text)
                run.font.name = 'Courier New'
                run.font.size = Pt(10)
            elif element.name == 'table':
                # 计算行列数
                rows = len(element.find_all('tr'))
                if rows > 0:
                    cols = max(len(row.find_all(['td', 'th'])) for row in element.find_all('tr'))

                    # 创建表格
                    table = doc.add_table(rows=rows, cols=cols)
                    table.style = 'Table Grid'

                    # 填充表格
                    for i, tr in enumerate(element.find_all('tr')):
                        cells = tr.find_all(['td', 'th'])
                        for j, cell in enumerate(cells):
                            if j < cols:  # 防止索引越界
                                table_cell = table.cell(i, j)
                                table_cell.text = cell.get_text().strip()
                                # 如果是表头
                                if cell.name == 'th':
                                    for paragraph in table_cell.paragraphs:
                                        for run in paragraph.runs:
                                            run.bold = True
            elif element.name == 'br':
                parent.add_run('\n')
            elif element.name == 'hr':
                doc.add_paragraph('_' * 50)
            elif element.name == 'blockquote':
                p = doc.add_paragraph()
                p.style = 'Quote'
                for child in element.children:
                    process_element(child, p)
            elif element.name == 'img':
                p = doc.add_paragraph()
                if 'src' in element.attrs:
                    src = element['src']
                    if src.startswith(('http://', 'https://')):
                        # 从URL下载图片
                        try:
                            response = requests.get(src)
                            image_data = BytesIO(response.content)
                            run = p.add_run()
                            run.add_picture(image_data)
                        except:
                            run = p.add_run(f"[图片: {src}]")
                    elif os.path.exists(src):
                        # 本地图片
                        run = p.add_run()
                        run.add_picture(src)
                    else:
                        run = p.add_run(f"[图片: {src}]")
            else:
                # 处理其他元素的子元素
                for child in element.children:
                    process_element(child, parent)

        # 处理整个文档
        for element in soup.children:
            process_element(element, doc)

    except Exception as e:
        doc.add_paragraph(f"Markdown转换失败: {str(e)}")
        doc.add_paragraph(markdown_text)  # 添加原始文本


def process_word_document(args: WriteToDocument) -> str:
    """处理Word文档，增强版"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
        from docx.oxml.ns import qn

        content_data = get_content_data(args)
        position_info = parse_position_word(args.position)
        font_mapping = parse_font_mapping(args.font_mapping)

        # 检查文档是否被占用
        doc_exists = os.path.exists(args.target_document)

        if doc_exists:
            try:
                # 尝试打开文档
                doc = Document(args.target_document)
            except Exception as e:
                error_msg = str(e).lower()
                if "被另一个进程占用" in error_msg or "access denied" in error_msg:
                    # 检查是否允许强制解锁
                    if hasattr(args, 'force_unlock') and args.force_unlock:
                        unlocked = force_unlock_document(args.target_document)
                        if unlocked:
                            # 重新尝试打开
                            try:
                                doc = Document(args.target_document)
                            except:
                                # 强制解锁失败，尝试只读模式
                                raise McpError(ErrorData(code=DOC_LOCKED, message=f"无法解除文档占用: {str(e)}"))
                        else:
                            # 解锁失败，尝试以只读模式打开
                            try:
                                # 创建临时副本
                                temp_path = tempfile.mktemp(suffix='.docx')
                                import shutil
                                shutil.copy2(args.target_document, temp_path)
                                doc = Document(temp_path)
                                # 设置标志，表明这是临时文件
                                is_temp_file = True
                            except:
                                raise McpError(
                                    ErrorData(code=DOC_LOCKED, message=f"文档被占用且无法以只读模式打开: {str(e)}"))
                    else:
                        # 未启用强制解锁，直接报错
                        raise McpError(ErrorData(code=DOC_LOCKED, message=f"文档被占用: {str(e)}"))
                elif args.password and ("密码" in str(e) or "password" in error_msg):
                    pass
                    # 尝试使用密码打开（需要win32com）
                    # try:
                    #     # 使用COM接口打开加密文档
                    #     pythoncom.CoInitialize()
                    #     word_app = win32com.client.Dispatch("Word.Application")
                    #     word_app.Visible = False
                    #
                    #     # 创建临时文件路径
                    #     temp_path = tempfile.mktemp(suffix='.docx')
                    #
                    #     # 打开原始文档并保存为未加密的临时文件
                    #     doc_com = word_app.Documents.Open(
                    #         os.path.abspath(args.target_document),
                    #         False,  # 确认转换
                    #         False,  # 只读
                    #         False,  # 添加到最近打开
                    #         args.password  # 密码
                    #     )
                    #
                    #     doc_com.SaveAs(os.path.abspath(temp_path), 16)  # 16 = wdFormatDocumentDefault
                    #     doc_com.Close()
                    #     word_app.Quit()
                    #
                    #     # 打开解密后的临时文档
                    #     doc = Document(temp_path)
                    #     # 设置标志，表明这是临时文件
                    #     is_temp_file = True
                    # except Exception as com_error:
                    #     raise McpError(
                    #         ErrorData(code=INVALID_PARAMS, message=f"文档密码错误或无法处理加密文档: {str(com_error)}"))
                else:
                    # 其他错误
                    raise ValueError(f"无法打开文档: {str(e)}")
        else:
            # 新建文档
            doc = Document()

        # 内容处理 - 根据内容类型
        if isinstance(content_data, bytes):  # 图片或二进制数据
            # 确定插入位置
            if position_info["position_type"] == "end":
                p = doc.add_paragraph()
            elif position_info["position_type"] == "start":
                p = doc.paragraphs[0] if doc.paragraphs else doc.add_paragraph()
            elif position_info["position_type"] == "paragraph":
                if position_info["paragraph_index"] < len(doc.paragraphs):
                    p = doc.paragraphs[position_info["paragraph_index"]]
                else:
                    p = doc.add_paragraph()
            elif position_info["position_type"] == "bookmark":
                # 处理书签位置
                bookmark_name = position_info["bookmark_name"]
                bookmark_found = False

                # XML操作查找书签
                for bookmark_start in doc.part.document.body.xpath("//w:bookmarkStart"):
                    if bookmark_start.get(qn("w:name")) == bookmark_name:
                        # 找到书签，在此创建段落
                        parent = bookmark_start.getparent()
                        if parent is not None:
                            p = doc.add_paragraph()
                            # 向父元素插入新段落
                            parent.insert(list(parent).index(bookmark_start) + 1, p._p)
                            bookmark_found = True
                            break

                if not bookmark_found:
                    p = doc.add_paragraph()
                    p.add_run(f"[书签 '{bookmark_name}' 未找到]")
            else:
                p = doc.add_paragraph()

            # 如果内容是多页PDF转换的图片列表
            if isinstance(content_data, list) and all(isinstance(item, bytes) for item in content_data):
                for i, image_data in enumerate(content_data):
                    # 为每个图像创建临时文件
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                        tmp.write(image_data)
                        tmp_name = tmp.name

                    # 插入图片
                    image_scale = parse_image_scale(args.image_scale)
                    if i > 0:
                        # 添加分页符
                        run = p.add_run()
                        run.add_break(WD_BREAK.PAGE)

                    run = p.add_run()

                    if "percent" in image_scale:
                        # 需要先获取图像尺寸
                        img = Image.open(tmp_name)
                        width, height = img.size
                        width_inches = width / 96  # 假设96 DPI
                        height_inches = height / 96

                        # 应用缩放
                        width_inches *= image_scale["percent"]
                        height_inches *= image_scale["percent"]

                        run.add_picture(tmp_name, width=Inches(width_inches), height=Inches(height_inches))
                    elif "width" in image_scale and "height" in image_scale:
                        width = Inches(image_scale["width"] / 96)  # convert px to inches
                        height = Inches(image_scale["height"] / 96)
                        run.add_picture(tmp_name, width=width, height=height)
                    else:
                        run.add_picture(tmp_name)

                    # 删除临时文件
                    os.unlink(tmp_name)
            else:
                # 单个图像
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                    tmp.write(content_data)
                    tmp_name = tmp.name

                # 插入图片
                image_scale = parse_image_scale(args.image_scale)
                run = p.add_run()

                if "percent" in image_scale:
                    # 需要先获取图像尺寸
                    img = Image.open(tmp_name)
                    width, height = img.size
                    width_inches = width / 96  # 假设96 DPI
                    height_inches = height / 96

                    # 应用缩放
                    width_inches *= image_scale["percent"]
                    height_inches *= image_scale["percent"]

                    run.add_picture(tmp_name, width=Inches(width_inches), height=Inches(height_inches))
                elif "width" in image_scale and "height" in image_scale:
                    width = Inches(image_scale["width"] / 96)  # convert px to inches
                    height = Inches(image_scale["height"] / 96)
                    run.add_picture(tmp_name, width=width, height=height)
                else:
                    run.add_picture(tmp_name)

                # 删除临时文件
                os.unlink(tmp_name)

        else:  # 文本或结构化数据
            # 检查插入模式
            insert_mode = args.insert_mode or "append"

            # 检查是否为Markdown格式
            is_markdown = False
            if args.content_type == "text" and (
                    args.content.startswith("#") or  # 标题
                    "```" in args.content or  # 代码块
                    "*" in args.content or  # 强调或列表
                    "- " in args.content  # 列表
            ):
                is_markdown = True

            # 如果是Markdown，使用特殊处理
            if is_markdown:
                # 如果是覆盖模式且不是在特定位置，则清空文档
                if insert_mode == "overwrite" and position_info["position_type"] in ["start", "end"]:
                    for i in range(len(doc.paragraphs) - 1, -1, -1):
                        p = doc.paragraphs[i]
                        p._element.getparent().remove(p._element)

                # 转换并插入Markdown
                markdown_to_docx(content_data, doc)

            # 如果是JSON格式的表格数据
            elif (args.content_type == "text" or args.content_type == "data") and (
                    isinstance(content_data, list) or
                    (isinstance(content_data, str) and content_data.strip().startswith(
                        "[") and content_data.strip().endswith("]"))
            ):
                table_data = content_data
                if isinstance(content_data, str):
                    try:
                        table_data = json.loads(content_data)
                    except:
                        # 如果解析失败，作为普通文本处理
                        table_data = None

                if isinstance(table_data, list) and all(isinstance(row, list) for row in table_data):
                    # 如果是覆盖模式且在特定位置有表格，则修改该表格
                    target_table = None
                    if insert_mode == "overwrite" and position_info["position_type"] == "paragraph":
                        para_idx = position_info["paragraph_index"]

                        # 查找段落后的表格
                        if para_idx < len(doc.paragraphs):
                            para = doc.paragraphs[para_idx]
                            para_idx_in_doc = -1

                            # 找到段落在文档中的位置
                            for i, p in enumerate(doc.paragraphs):
                                if p._p == para._p:
                                    para_idx_in_doc = i
                                    break

                            if para_idx_in_doc >= 0:
                                # 检查段落后是否紧跟着一个表格
                                element = para._p.getnext()
                                if element is not None and element.tag.endswith('tbl'):
                                    # 找到对应的表格对象
                                    for tbl in doc.tables:
                                        if tbl._tbl == element:
                                            target_table = tbl
                                            break

                    # 创建或修改表格
                    rows = len(table_data)
                    cols = max(len(row) for row in table_data)

                    if target_table and insert_mode == "overwrite":
                        # 调整现有表格大小
                        current_rows = len(target_table.rows)
                        current_cols = len(target_table.columns)

                        # 添加或删除行
                        if current_rows < rows:
                            # 添加行
                            for _ in range(rows - current_rows):
                                target_table.add_row()
                        elif current_rows > rows:
                            # 删除行（从底部开始）
                            for _ in range(current_rows - rows):
                                tr = target_table._tbl.tr_lst[-1]
                                target_table._tbl.remove(tr)

                        # 添加列（无法直接删除列，只能添加）
                        if current_cols < cols:
                            for row in target_table.rows:
                                for _ in range(cols - current_cols):
                                    cell = row.cells[-1]._tc.add_successor('w:tc')
                                    cell.text = ''

                        # 填充数据
                        for i, row in enumerate(table_data):
                            for j, cell_data in enumerate(row):
                                if j < cols:  # 确保不超出列范围
                                    if j >= current_cols:
                                        # 跳过新添加的超出原始列范围的列
                                        continue
                                    cell = target_table.cell(i, j)
                                    cell.text = str(cell_data)
                    else:
                        # 创建新表格
                        table = doc.add_table(rows=rows, cols=cols)

                        # 设置表格样式
                        if args.table_style:
                            try:
                                table.style = args.table_style
                            except:
                                pass  # 忽略不支持的样式

                        # 填充表格数据
                        for i, row in enumerate(table_data):
                            for j, cell_data in enumerate(row):
                                if j < cols:  # 确保不超出列范围
                                    cell = table.cell(i, j)
                                    cell.text = str(cell_data)
                else:
                    # 不是有效的表格数据，作为普通文本添加
                    if position_info["position_type"] == "end":
                        p = doc.add_paragraph(str(content_data))
                    elif position_info["position_type"] == "start":
                        p = doc.paragraphs[0] if doc.paragraphs else doc.add_paragraph()
                        if insert_mode == "overwrite":
                            p.text = str(content_data)
                        else:
                            p.insert_paragraph_before(str(content_data))
                    elif position_info["position_type"] == "paragraph":
                        if position_info["paragraph_index"] < len(doc.paragraphs):
                            p = doc.paragraphs[position_info["paragraph_index"]]
                            if insert_mode == "overwrite":
                                p.text = str(content_data)
                            else:
                                p.add_run(str(content_data))
                        else:
                            p = doc.add_paragraph(str(content_data))
                    elif position_info["position_type"] == "bookmark":
                        # 与前面图片处理类似，重复查找书签的代码
                        bookmark_name = position_info["bookmark_name"]
                        bookmark_found = False

                        for bookmark_start in doc.part.document.body.xpath("//w:bookmarkStart"):
                            if bookmark_start.get(qn("w:name")) == bookmark_name:
                                parent = bookmark_start.getparent()
                                if parent is not None:
                                    p = doc.add_paragraph(str(content_data) if insert_mode == "overwrite" else "")
                                    parent.insert(list(parent).index(bookmark_start) + 1, p._p)
                                    bookmark_found = True
                                    break

                        if not bookmark_found:
                            p = doc.add_paragraph()
                            p.add_run(f"[书签 '{bookmark_name}' 未找到] ")
                            p.add_run(str(content_data))
            else:
                # 普通文本
                if position_info["position_type"] == "end":
                    p = doc.add_paragraph(content_data)
                elif position_info["position_type"] == "start":
                    p = doc.paragraphs[0] if doc.paragraphs else doc.add_paragraph()
                    if insert_mode == "overwrite":
                        p.text = content_data
                    else:
                        p.insert_paragraph_before(content_data)
                elif position_info["position_type"] == "paragraph":
                    if position_info["paragraph_index"] < len(doc.paragraphs):
                        p = doc.paragraphs[position_info["paragraph_index"]]
                        if insert_mode == "overwrite":
                            p.text = content_data
                        else:
                            p.add_run(content_data)
                    else:
                        p = doc.add_paragraph(content_data)
                elif position_info["position_type"] == "bookmark":
                    bookmark_name = position_info["bookmark_name"]
                    bookmark_found = False

                    for bookmark_start in doc.part.document.body.xpath("//w:bookmarkStart"):
                        if bookmark_start.get(qn("w:name")) == bookmark_name:
                            parent = bookmark_start.getparent()
                            if parent is not None:
                                p = doc.add_paragraph(content_data if insert_mode == "overwrite" else "")
                                parent.insert(list(parent).index(bookmark_start) + 1, p._p)
                                bookmark_found = True
                                break

                    if not bookmark_found:
                        p = doc.add_paragraph()
                        p.add_run(f"[书签 '{bookmark_name}' 未找到] ")
                        p.add_run(content_data)

        # 应用字体映射
        if font_mapping:
            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    if run.font.name in font_mapping:
                        run.font.name = font_mapping[run.font.name]

        # 处理版本兼容性
        if args.compatibility_mode and args.compatibility_mode.lower() == "97-2003":
            # 保存为旧版本格式
            old_format = True
        else:
            old_format = False

        # 保存文档
        output_path = get_output_path(args.target_document, args.overwrite_save == 'true', args.output_format)

        # 确保输出目录存在
        ensure_dir_exists(output_path)

        # 如果需要以旧格式保存
        if old_format and output_path.endswith('.docx'):
            output_path = output_path.replace('.docx', '.doc')

        doc.save(output_path)

        return f"Word文档处理成功，内容已插入到 {output_path}"

    except ImportError as e:
        if "docx" in str(e):
            raise ImportError("python-docx")
        elif "com" in str(e) or "win32" in str(e):
            raise ImportError("pywin32")
        elif "markdown" in str(e):
            raise ImportError("markdown")
        elif "beautifulsoup4" in str(e) or "bs4" in str(e):
            raise ImportError("beautifulsoup4")
        else:
            raise ImportError(str(e))
    except Exception as e:
        raise Exception(f"处理Word文档时出错: {str(e)}")




def process_powerpoint_document(args: WriteToDocument) -> str:
    """处理PowerPoint文档"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        content_data = get_content_data(args)
        position_info = parse_position_powerpoint(args.position)
        
        # 打开现有演示文稿或创建新演示文稿
        prs = None
        if os.path.exists(args.target_document):
            try:
                prs = Presentation(args.target_document)
            except Exception as e:
                if args.password:
                    raise ValueError(f"无法打开受密码保护的文档: {str(e)}")
                else:
                    raise ValueError(f"无法打开文档: {str(e)}")
        else:
            prs = Presentation()
        
        # 确定要操作的幻灯片
        slide = None
        if position_info["slide"] == "end":
            # 在末尾添加新幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
        elif position_info["slide"] == "new":
            # 添加新幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白布局
        else:
            # 使用指定的幻灯片
            if 0 <= position_info["slide"] < len(prs.slides):
                slide = prs.slides[position_info["slide"]]
            else:
                # 如果指定的幻灯片不存在，创建新幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # 根据内容类型处理
        if isinstance(content_data, bytes):  # 二进制数据(图片)
            # 创建临时文件保存图片
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                tmp.write(content_data)
                tmp_name = tmp.name
            
            # 插入图片
            image_scale = parse_image_scale(args.image_scale)
            
            if "percent" in image_scale:
                # 按百分比缩放，需要先获取原图尺寸
                width = Inches(5) * image_scale["percent"]  # 假设原始尺寸
                height = Inches(4) * image_scale["percent"]
                slide.shapes.add_picture(tmp_name, Inches(1), Inches(1), width=width, height=height)
            elif "width" in image_scale and "height" in image_scale:
                # 指定尺寸
                width = Inches(image_scale["width"] / 96)  # convert px to inches
                height = Inches(image_scale["height"] / 96)
                slide.shapes.add_picture(tmp_name, Inches(1), Inches(1), width=width, height=height)
            else:
                # 原始尺寸
                slide.shapes.add_picture(tmp_name, Inches(1), Inches(1))
            
            # 删除临时文件
            os.unlink(tmp_name)
            
        else:  # 文本内容
            # 检查是否为JSON表格数据
            if args.content_type == "text" and content_data.strip().startswith("[") and content_data.strip().endswith("]"):
                try:
                    table_data = json.loads(content_data)
                    if isinstance(table_data, list) and all(isinstance(row, list) for row in table_data):
                        # 创建表格
                        rows = len(table_data)
                        cols = max(len(row) for row in table_data)
                        
                        # 指定表格大小和位置
                        left = Inches(1)
                        top = Inches(1)
                        width = Inches(8)
                        height = Inches(rows * 0.5)
                        
                        # 添加表格
                        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                        
                        # 填充表格数据
                        for i, row in enumerate(table_data):
                            for j, cell_data in enumerate(row):
                                if j < cols:  # 确保不超出列范围
                                    table.cell(i, j).text = str(cell_data)
                    else:
                        # 添加为普通文本
                        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                        text_frame = text_box.text_frame
                        text_frame.text = content_data
                except json.JSONDecodeError:
                    # 不是有效的JSON，作为普通文本添加
                    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                    text_frame = text_box.text_frame
                    text_frame.text = content_data
            else:
                # 普通文本，添加文本框
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                text_frame = text_box.text_frame
                
                # 处理多行文本
                lines = content_data.split('\n')
                text_frame.text = lines[0] if lines else ""
                
                # 添加额外的行
                for line in lines[1:]:
                    p = text_frame.add_paragraph()
                    p.text = line
        
        # 保存演示文稿
        output_path = get_output_path(args.target_document, args.overwrite_save == 'true', args.output_format)
        prs.save(output_path)
        
        return f"PowerPoint文档处理成功，内容已插入到 {output_path}"
    
    except ImportError:
        raise ImportError("python-pptx")
    except Exception as e:
        raise Exception(f"处理PowerPoint文档时出错: {str(e)}")

async def process_batch_documents(args: BatchWriteToDocument) -> str:
    """批量处理文档"""
    try:
        if not os.path.exists(args.input_list):
            return f"错误：输入列表文件 {args.input_list} 不存在"
        
        with open(args.input_list, 'r', encoding='utf-8') as f:
            documents = json.load(f)
        
        results = []
        for i, doc in enumerate(documents):
            try:
                doc_args = WriteToDocument(**doc)
                result = await process_document(doc_args)
                results.append(f"[{i+1}/{len(documents)}] {doc_args.target_document}: 成功")
            except Exception as e:
                results.append(f"[{i+1}/{len(documents)}] {doc.get('target_document', 'Unknown')}: 失败 - {str(e)}")
        
        return "\n".join(results)
    except Exception as e:
        return f"批量处理文档时出错: {str(e)}"

async def revert_document(args: RevertDocument) -> str:
    """恢复文档到指定版本"""
    try:
        base_name = os.path.splitext(args.target_document)[0]
        ext = os.path.splitext(args.target_document)[1]
        backup_path = f"{base_name}_version{args.version}.bak"
        
        if not os.path.exists(backup_path):
            return f"错误：备份文件 {backup_path} 不存在"
        
        import shutil
        # 先备份当前版本
        current_backup = create_backup(args.target_document)
        # 然后恢复指定版本
        shutil.copy2(backup_path, args.target_document)
        
        return f"文档已恢复到版本 {args.version}，当前版本已备份到 {current_backup}"
    except Exception as e:
        return f"恢复文档时出错: {str(e)}"

def force_unlock_document(file_path: str) -> bool:
    """尝试强制解锁被占用的文档（Windows系统专用）"""
    try:
        if os.name != 'nt':  # 只在Windows系统下尝试
            return False
        
        import subprocess
        # 使用Windows系统命令查找占用文件的进程
        result = subprocess.run(
            ["handle", file_path], 
            capture_output=True, 
            text=True, 
            shell=True
        )
        
        # 解析输出查找进程ID
        lines = result.stdout.splitlines()
        pids = []
        for line in lines:
            if file_path in line:
                parts = line.split()
                if len(parts) > 2:
                    pid_part = parts[2]
                    if pid_part.startswith('pid:'):
                        pid = pid_part[4:]
                        pids.append(pid)
        
        # 结束占用进程
        for pid in pids:
            subprocess.run(["taskkill", "/F", "/PID", pid], shell=True)
        
        return len(pids) > 0
    except Exception:
        return False

async def serve() -> None:
    server = Server("mcp-document-writer")

    @server.list_tools()
    async def list_tools() -> list[Tool]:
        return [
            Tool(name="write_to_document", description="将内容写入Office文档", inputSchema=WriteToDocument.model_json_schema()),
            Tool(name="batch_write_to_document", description="批量处理多个文档", inputSchema=BatchWriteToDocument.model_json_schema()),
            Tool(name="revert_document", description="恢复文档到指定版本", inputSchema=RevertDocument.model_json_schema()),
        ]

    @server.call_tool()
    async def call_tool(name: str, arguments: dict) -> list[TextContent]:
        try:
            if name == "write_to_document":
                args = WriteToDocument(**arguments)
                result = await process_document(args)
            elif name == "batch_write_to_document":
                args = BatchWriteToDocument(**arguments)
                result = await process_batch_documents(args)
            elif name == "revert_document":
                args = RevertDocument(**arguments)
                result = await revert_document(args)
            else:
                raise ValueError(f"未知的工具名称: {name}")

            return [TextContent(type="text", text=result)]
        except McpError as e:
            raise e
        except Exception as e:
            raise McpError(ErrorData(code=INTERNAL_ERROR, message=str(e)))

    @server.list_prompts()
    async def list_prompts() -> list[Prompt]:
        return []

    @server.get_prompt()
    async def get_prompt(name: str, arguments: dict | None) -> GetPromptResult:
        raise McpError(ErrorData(code=INVALID_PARAMS, message="不支持的操作"))

    options = server.create_initialization_options()
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, options, raise_exceptions=True)

def get_clipboard_content():
    """获取剪贴板内容，支持文本/图片/文件"""
    try:
        import pyperclip
        import PIL.ImageGrab
        
        # 尝试获取图片
        image = PIL.ImageGrab.grabclipboard()
        if image:
            # 返回图片对象
            img_bytes = BytesIO()
            image.save(img_bytes, format='PNG')
            return img_bytes.getvalue()
        
        # 尝试获取文件列表
        if hasattr(PIL.ImageGrab, 'grabclipboard') and callable(PIL.ImageGrab.grabclipboard):
            files = PIL.ImageGrab.grabclipboard()
            if isinstance(files, list) and len(files) > 0 and all(isinstance(f, str) for f in files):
                return files[0]  # 返回第一个文件路径
        
        # 获取文本
        text = pyperclip.paste()
        if text:
            return text
            
        return None
    except Exception as e:
        raise Exception(f"获取剪贴板内容失败: {str(e)}")


def get_http_content(url: str) -> Any:
    """从HTTP端点获取内容"""
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        
        # 检查内容类型
        content_type = response.headers.get('Content-Type', '')
        
        if 'application/json' in content_type:
            return response.json()
        elif 'text/' in content_type:
            return response.text
        elif 'image/' in content_type:
            return response.content
        else:
            return response.content
    except Exception as e:
        raise Exception(f"从HTTP端点获取内容失败: {str(e)}")


def check_image_resolution(image_data: bytes) -> bool:
    """检查图片分辨率是否超过限制"""
    try:
        from PIL import Image
        from io import BytesIO
        
        img = Image.open(BytesIO(image_data))
        width, height = img.size
        
        if width > 10000 or height > 10000:
            return False
        return True
    except Exception:
        return True  # 如果无法检查，默认允许


def log_table_conflict(target_document: str, sheet_name: str, cell_address: str, original_value: Any, new_value: Any):
    """记录表格合并冲突"""
    from datetime import datetime
    log_file = f"{os.path.splitext(target_document)[0]}_conflict.log"
    
    with open(log_file, 'a', encoding='utf-8') as f:
        f.write(f"时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"文档: {target_document}\n")
        f.write(f"位置: {sheet_name}!{cell_address}\n")
        f.write(f"原值: {original_value}\n")
        f.write(f"新值: {new_value}\n")
        f.write("-" * 50 + "\n")


def check_and_strip_vba(file_path: str) -> bool:
    """检查并剥离VBA代码，返回是否含有VBA"""
    try:
        if file_path.endswith(('.docm', '.xlsm', '.pptm')):
            # 这里需要使用适当的库检查VBA
            # 例如使用oletools库
            import oletools.olevba
            vba_parser = oletools.olevba.VBA_Parser(file_path)
            has_vba = vba_parser.detect_vba_macros()
            
            if has_vba:
                # 创建无宏版本
                base_name = os.path.splitext(file_path)[0]
                ext = os.path.splitext(file_path)[1]
                new_ext = ext.replace('m', 'x')  # 将.docm改为.docx等
                new_path = f"{base_name}_no_macros{new_ext}"
                
                # 这里需要实现文档转换逻辑
                # ...
                
                return True
            return False
    except ImportError:
        # 如果没有oletools库，返回警告
        print("警告: 无法检查VBA代码，oletools库未安装")
        return False
    except Exception as e:
        print(f"检查VBA代码时出错: {str(e)}")
        return False


def main():
    print("Document Writer Server 正在启动...")
    import asyncio
    asyncio.run(serve())

if __name__ == '__main__':
    main()
